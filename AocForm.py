import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import tempfile
from datetime import datetime
import requests
from io import BytesIO
from PIL import Image

# Streamlit UI
st.title("Agent of Change (AOC) Form")

# Upload Header Image
#header_image = st.file_uploader("Upload Header Image", type=["jpg", "png"])
header_image = "https://github.com/Febrinaads/Febrina-Dyah-Sukmawati/blob/main/Header%20Format.png?raw=true"

headers = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
}

try:
    response = requests.get(header_image, headers=headers, timeout=10)
    response.raise_for_status()  # Check if request was successful

    image = Image.open(BytesIO(response.content))  # Open the image from bytes
    image.save("header_image.png")  # Save the image locally
    print("Image successfully downloaded and saved!")

except requests.exceptions.RequestException as e:
    print(f"Failed to load the header image: {e}")

#header_image = "https://raw.githubusercontent.com/Febrinaads/Febrina-Dyah-Sukmawati/main/Header%20Format.png?raw=true" #NEED PROPER LINK
#response = requests.get(header_image)

# Input Fields
st.header("Form Penjelasan Dokumentasi")
teks0 = st.text_input("Jenis Kegiatan: ")

st.header("Data Narasumber")
teks1 = st.text_input("Nama Lengkap: ")
teks2 = st.text_input("NIK: ")
teks3 = st.text_input("Jabatan: ")
teks4 = st.text_input("Email: ")
teks5 = st.text_input("Holding/Sub Holding: ")
teks6 = st.text_input("Perusahaan: ")
teks7 = st.text_input("Direktorat: ")

gambar_list = st.file_uploader("Upload Images (Max 3)", type=["jpg", "png"], accept_multiple_files=True)
if len(gambar_list) > 3:
    st.warning("You can only upload up to 3 images.")
    gambar_list = gambar_list[:3]

st.header("Jadwal Pelaksanaan")
teks8 = st.date_input("Hari/Tanggal: ", value=datetime.today().date())
teks9 = st.time_input("Waktu Mulai", value=datetime.strptime("08:30", "%H:%M").time())
teks9b = st.time_input("Waktu Selesai", value=datetime.strptime("12:00", "%H:%M").time())

st.header("Deskripsi Kegiatan")
MAX_CHAR = 700
teks10 = st.text_area("AKHLAK Moment (Max 700 characters):", max_chars=MAX_CHAR)
remaining_chars = MAX_CHAR - len(teks10)
if len(teks10) > MAX_CHAR:
    st.warning(f"Character limit exceeded! Only the first 700 characters will be saved.")
    teks10 = teks10[:MAX_CHAR]  # Truncate text if needed
    
if st.button("Save to PowerPoint"):
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # Set width for 16:9 ratio
    prs.slide_height = Inches(7.5)   # Set height for 16:9 ratio
    
    slide_layout = prs.slide_layouts[5]  # Blank Layout
    slide = prs.slides.add_slide(slide_layout)
    
    if response.status_code == 200:
        try:
            image = Image.open(BytesIO(response.content))  # Validate image
            temp_header_path = "header_image.png"
            image.save(temp_header_path)  # Save it temporarily
            
            # Use slide AFTER it is defined
            slide.shapes.add_picture(temp_header_path, Inches(0), Inches(0), width=Inches(13.33))
        except Exception as e:
            st.error(f"Error processing image: {e}")
    else:
        st.error("Failed to load the header image. Please check the URL.")
    
    # Title
    title = slide.shapes.title
    title.text = "FORM PENJELASAN DOKUMENTASI"
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Inches(0.33)
    title.left = Inches(2)
    title.top = Inches(0.25)
    title.width = Inches(10)
    title.height = Inches(0.25)
    
    # Table for Jenis Kegiatan
    table0 = slide.shapes.add_table(1, 2, Inches(0.6), Inches(0.75), Inches(6), Inches(0.5)).table
    table0.columns[0].width = Inches(3)
    table0.columns[1].width = Inches(3)
    table0.cell(0, 0).text = "Jenis Kegiatan:"
    table0.cell(0, 1).text = teks0

    #create Data Narasumber Textbox
    text_box = slide.shapes.add_textbox(Inches(2.5), Inches(1.30), Inches(2), Inches(0.5))
    text_frame = text_box.text_frame
    text_frame.text = "Data Narasumber"
    
    # Table for Data Narasumber
    table1 = slide.shapes.add_table(7, 2, Inches(0.6), Inches(1.75), Inches(6), Inches(0.5)).table
    table1.columns[0].width = Inches(3)
    table1.columns[1].width = Inches(3)
    labels = ["Nama Lengkap", "NIK", "Jabatan", "Email", "Holding/Sub Holding", "Perusahaan", "Direktorat"]
    values = [teks1, teks2, teks3, teks4, teks5, teks6, teks7]
    for i, (label, value) in enumerate(zip(labels, values)):
        table1.cell(i, 0).text = label
        table1.cell(i, 1).text = value

    #create Jadwal Pelaksanaan Textbox
    text_box = slide.shapes.add_textbox(Inches(8.5), Inches(0.75), Inches(2), Inches(0.5))
    text_frame = text_box.text_frame
    text_frame.text = "Jadwal Pelaksanaan"
    
    # Table for Jadwal Pelaksanaan
    table2 = slide.shapes.add_table(2, 2, Inches(6.8), Inches(1.30), Inches(3), Inches(0.5)).table
    table2.columns[0].width = Inches(3)
    table2.columns[1].width = Inches(3)
    table2.cell(0, 0).text = "Hari/Tanggal"
    table2.cell(0, 1).text = teks8.strftime('%d-%m-%Y')
    table2.cell(1, 0).text = "Waktu (Start – End)"
    table2.cell(1, 1).text = f"{teks9.strftime('%H:%M')} - {teks9b.strftime('%H:%M')} WIB"
    
    # Table for Deskripsi Kegiatan
    table3 = slide.shapes.add_table(2, 1, Inches(6.8), Inches(2.25), Inches(6), Inches(0.5)).table
    table3.columns[0].width = Inches(6)
    table3.cell(0, 0).text = "Deskripsi Kegiatan"
    table3.cell(1, 0).text = teks10

    #create Dokumentasi Textbox
    text_box = slide.shapes.add_textbox(Inches(2.5), Inches(4.88), Inches(2), Inches(0.5))
    text_frame = text_box.text_frame
    text_frame.text = "Dokumentasi"
    
    # Table for Images
    table4 = slide.shapes.add_table(1, 1, Inches(0.6), Inches(5.25), Inches(6), Inches(2)).table
    if gambar_list:
        x_offset = Inches(1.25)
        y_offset = Inches(5.3)
        max_height = Inches(1.75)
        img_spacing = Inches(1.5)
        for gambar in gambar_list:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_file:
                tmp_file.write(gambar.read())
                temp_image_path = tmp_file.name
            slide.shapes.add_picture(temp_image_path, x_offset, y_offset, height=max_height)
            x_offset += img_spacing
    
    # Save file
    pptx_path = "AoC_Form.pptx"
    prs.save(pptx_path)
    st.success("PowerPoint file created successfully!")
    with open(pptx_path, "rb") as file:
        st.download_button("Download PowerPoint", file, file_name="AoC_Form.pptx")